''' 文本提取模块
1）用于从不同格式的文档中提取文本信息，支持docx、pdf、txt、excel、pptx等格式
2）用于从文本中提取实体信息，支持正则表达式规则
'''
from tqdm import tqdm
import os
import threading
import multiprocessing
from typing import Optional, Union, List, Tuple, Literal, Dict, Any
import shutil
import json
import re
import numpy as np
from torch4keras.snippets import TryExcept, safe_import, log_warn, log_error
with safe_import():
    import pdfplumber
    from pdfplumber.page import Page as PdfPlumberPage
with safe_import():
    import fitz  # PyMuPDF


def extract_entity_by_rule( 
    text: str,
    pattern: Optional[str] = None,
    label: Optional[str] = None,
    start: int = 0,
    end: int = -1,
    dotall: bool = True,
    replace_pattern: Optional[Union[str, List[str]]] = None,
    extract_pattern: Optional[Union[str, List[str]]] = None,
    min_entity_len: Optional[int] = None,
    max_entity_len: Optional[int] = None,
    ignore_chars: Optional[Union[str, List[str], Tuple[str, ...]]] = None,
    exist_subword: Optional[Union[List[str], str, Tuple[str, ...]]] = None,
    noexist_subword: Optional[Union[List[str], str, Tuple[str, ...]]] = None,
    prefix_exist_subword: Optional[List[Tuple[str, int]]] = None,
    prefix_noexist_subword: Optional[List[Tuple[str, int]]] = None,
    postfix_exist_subword: Optional[List[Tuple[str, int]]] = None,
    postfix_noexist_subword: Optional[List[Tuple[str, int]]] = None,
    **kwargs,  # 保留扩展性
 ) -> List[Dict[str, Any]]:  # 明确返回类型
    """按照预设的正则规则从文本中提取实体。
 
    Args:
        text: 待提取的文本。
        pattern: 用于提取的正则表达式模式。
        label: 提取结果的标签（如实体类型）。
        start: 文本截取的起始位置（闭区间）。
        end: 文本截取的结束位置（闭区间）。
        dotall: 正则匹配时是否让 `.` 匹配换行符（`re.DOTALL`）。
        replace_pattern: 对正则匹配结果的二次修正（如去除前后缀），支持字符串或列表。
        extract_pattern: 从匹配结果中提取部分内容（如子组），支持字符串或列表。
        min_entity_len: 实体最小长度（低于此长度则丢弃）。
        max_entity_len: 实体最大长度（超过此长度则丢弃）。
        ignore_chars: 文本中需要忽略的字符（如空格、换行符）。
        exist_subword: 实体必须包含的子词（字符串或列表）。
        noexist_subword: 实体必须不包含的子词（字符串或列表）。
        prefix_exist_subword: 实体前缀必须包含的子词及其距离，格式为 `[('subword', distance), ...]`。
        prefix_noexist_subword: 实体前缀必须不包含的子词及其距离。
        postfix_exist_subword: 实体后缀必须包含的子词及其距离。
        postfix_noexist_subword: 实体后缀必须不包含的子词及其距离。
        **kwargs: 其他扩展参数。
 
    Returns:
        列表，每个元素是字典，包含提取的实体信息，例如：
        ```python
        [
            {
                "entity": "中国工商银行",  # 清理后的实体文本
                "raw_entity": "甲方：中国工商银行 乙方",  # 原始匹配文本
                "start": 3,  # 在原文本中的起始位置
                "end": 10,   # 在原文本中的结束位置
                "label": "甲方"  # 实体标签
            }
        ]
        ```
 
    Examples:
        >>> text = "甲方：中国工商银行 乙方：中国农业银行 注册地址：上海市世纪大道1379号"
        >>> config = {
        ...     "pattern": "甲方(:|：)(.*?)乙方",
        ...     "label": "甲方",
        ...     "replace_pattern": ["^甲方(:|：)", "乙方$"]
        ... }
        >>> res = extract_entity_by_rule(text, **config)
        >>> print(res)
        [{'entity': '中国工商银行', 'raw_entity': '甲方：中国工商银行 乙方', 'start': 3, 'end': 10, 'label': '甲方'}]
    """
    def adjust_start_end(entity, new_entity, start):
        if new_entity in entity:
            start += entity.index(new_entity)
            end = start + len(new_entity)
            return new_entity, start, end
        else:
            log_warn(f'{new_entity} <------- not in -------> {entity}')
            return entity, start, start + len(entity)

    def preprocess_text(text):
        # 替换空格和换行，并记录位置
        offset_map = []
        processed_text = []
        i = 0  # 映射是processed_text中的未知以及对应的offset
        for char in text:
            if char in ignore_chars:
                offset_map.append((i, 1))
            else:
                processed_text.append(char)
                i += 1
        return ''.join(processed_text), offset_map

    def restore_position(start, end, offset_map):
        # 根据offset_map恢复原始位置
        original_start = start
        original_end = end
        for pos, offset in offset_map:
            if pos <= start:
                original_start += offset
            if pos < end:  # 这里不包含等号
                original_end += offset
        return original_start, original_end

    # 截取一下
    if start != 0:
        text = text[start:]
    if end != -1:
        text = text[:end]

    # 预处理文本
    if ignore_chars is not None:
        processed_text, offset_map = preprocess_text(text)
    else:
        processed_text, offset_map = text, []

    # 中间.*可以匹配换行符
    iters = re.finditer(pattern, processed_text, re.DOTALL if dotall else 0)

    result = []
    for iter in iters:
        entity = raw_entity = iter.group()
        start, end = iter.start(), iter.end()
        original_raw_start, original_raw_end = restore_position(start, end, offset_map)

        # 提取的pattern
        if extract_pattern is not None:
            if isinstance(extract_pattern, str):
                extract_pattern = [extract_pattern]
            for pat in extract_pattern:
                if re.search(pat, entity):
                    new_entity = next(re.finditer(pat, entity)).group()
                    entity, start, end = adjust_start_end(entity, new_entity, start)

        # 删除的pattern
        if replace_pattern is not None:
            if isinstance(replace_pattern, str):
                replace_pattern = [replace_pattern]
            for rep_pat in replace_pattern:
                if re.search(rep_pat, entity):
                    new_entity = re.sub(rep_pat, '', entity)
                    entity, start, end = adjust_start_end(entity, new_entity, start)

        # 太短
        if (min_entity_len is not None) and (len(entity) <= min_entity_len):
            continue
        
        # 超长
        if (max_entity_len is not None) and (len(entity) >= max_entity_len):
            continue
        
        # exist_subword: 必须存在的subword
        if exist_subword is not None:
            if isinstance(exist_subword, str) and (not re.search(exist_subword, entity)):
                continue
            elif isinstance(exist_subword, (tuple, list)):
                continue_tag = False
                for item in exist_subword:
                    if not re.search(item, entity):
                        continue_tag = True
                        break
                if continue_tag:
                    continue

        # noexist_subword: 必须不存在的subword
        if noexist_subword is not None:
            if isinstance(noexist_subword, str) and re.search(noexist_subword, entity):
                continue
            elif isinstance(noexist_subword, (tuple, list)):
                continue_tag = False
                for item in noexist_subword:
                    if re.search(item, entity):
                        continue_tag = True
                        break
                if continue_tag:
                    continue
        
        # prefix_exist_subword: prefix中必须存在的subword
        if prefix_exist_subword is not None:
            assert isinstance(prefix_exist_subword, (tuple, list)), 'prefix_exist_subword only accept tuple/list format'
            prefix_exist_subword = [prefix_exist_subword] if isinstance(prefix_exist_subword[0], str) else prefix_exist_subword
            continue_tag = False
            for item, offset in prefix_exist_subword:
                if not re.search(item, processed_text[max(0, start-offset):start]):
                    continue_tag = True
                    break
            if continue_tag:
                continue

        # prefix_noexist_subword: prefix中必须不存在的subword
        if prefix_noexist_subword is not None:
            assert isinstance(prefix_noexist_subword, (tuple, list)), 'prefix_noexist_subword only accept tuple/list format'
            prefix_noexist_subword = [prefix_noexist_subword] if isinstance(prefix_noexist_subword[0], str) else prefix_noexist_subword
            continue_tag = False
            for item, offset in prefix_noexist_subword:
                if re.search(item, processed_text[max(0, start-offset):start]):
                    continue_tag = True
                    break
            if continue_tag:
                continue
        
        # postfix_exist_subword: postfix中必须存在的subword
        if postfix_exist_subword is not None:
            assert isinstance(postfix_exist_subword, (tuple, list)), 'postfix_exist_subword only accept tuple/list format'
            postfix_exist_subword = [postfix_exist_subword] if isinstance(postfix_exist_subword[0], str) else postfix_exist_subword
            continue_tag = False
            for item, offset in postfix_exist_subword:
                if not re.search(item, processed_text[end:end+offset]):
                    continue_tag = True
                    break
            if continue_tag:
                continue

        # postfix_noexist_subword: postfix中必须不存在的subword
        if postfix_noexist_subword is not None:
            assert isinstance(postfix_noexist_subword, (tuple, list)), 'postfix_noexist_subword only accept tuple/list format'
            postfix_noexist_subword = [postfix_noexist_subword] if isinstance(postfix_noexist_subword[0], str) else postfix_noexist_subword
            continue_tag = False
            for item, offset in postfix_noexist_subword:
                if re.search(item, processed_text[end:end+offset]):
                    continue_tag = True
                    break
            if continue_tag:
                continue
        
        # 恢复原始位置
        original_start, original_end = restore_position(start, end, offset_map)
        
        # 从原始文本中提取entity和raw_entity
        entity = text[original_start:original_end]
        raw_entity = text[original_raw_start:original_raw_end]
        if ignore_chars is not None:
            assert re.sub('|'.join(ignore_chars), '', entity) == processed_text[start:end]

        result.append(
            {
                'entity': entity, 
                'raw_entity': raw_entity,
                'start': original_start, 
                'end': original_end, 
                'label': label,
                **kwargs
            })
    return result


def extract_docx(src_docx:str):
    '''提取docx文本'''
    from docx import Document   # 导入相关库

    doc = Document(src_docx)
    full_text = []
    for p in tqdm(doc.paragraphs, desc='WORD extracting:', ncols=80):
        if p.text != '':
            full_text.append(p.text)
    full_text = '\n'.join(full_text)

    return full_text, doc


def extract_pdf(pdf_path:str, cache_dir:str=None, page_start:int=0, page_end:int=-1,
                method:Literal['pdfplumber', 'fitz']='pdfplumber', raise_error:bool=True, 
                replace:bool=False, num_processes:int=None, num_threads:int=None) -> List[str]:
    '''解析pdf文件
    测试发现fitz解析要比pdfplumber快很多倍，尤其是大文件

    :param pdf_path: pdf文件路径
    :param cache_dir: 缓存路径
    :param page_start: 开始页码
    :param page_end: 结束页码
    :param method: 解析方法, pdfplumber, fitz
    :param replace: 是否替换pdf
    '''
    if replace and cache_dir is not None:
        # 如果存在，删除里面所有文件夹
        shutil.rmtree(cache_dir, ignore_errors=True)
        
    # 如果不存在，则创建文件夹
    if cache_dir and not os.path.exists(cache_dir):
        os.makedirs(cache_dir, exist_ok=True)
    
    res = None
    with TryExcept(default=[], reraise=raise_error):
        if num_processes is not None and num_processes > 1:
            # 多进程
            res = extract_pdf_multiprocess(pdf_path, cache_dir, page_start, page_end, num_processes=num_processes, method=method)
        elif num_threads is not None and num_threads > 1:
            # 多线程
            res = extract_pdf_multithread(pdf_path, cache_dir, page_start, page_end, num_threads=num_threads, method=method)
        else:
            # 单进程
            res = extract_pdf_single_process(pdf_path, cache_dir, page_start, page_end, method=method)
    return res


class PDFParser:
    '''定义解析器接口'''
    @staticmethod
    def open(pdf_path: str):
        raise NotImplementedError
    
    @staticmethod
    def get_page_text(page: Any) -> str:
        raise NotImplementedError
    
    @staticmethod
    def get_page_number(page: Any) -> int:
        raise NotImplementedError
    
    @staticmethod
    def get_pages(pdf: Any) -> int:
        raise NotImplementedError


class PDFPlumberParser(PDFParser):
    ''' pdfplumber解析器实现 '''
    @staticmethod
    def open(pdf_path: str):
        return pdfplumber.open(pdf_path)
    
    @staticmethod
    def get_page_text(page: PdfPlumberPage) -> str:
        return page.extract_text() or ""
    
    @staticmethod
    def get_page_number(page: PdfPlumberPage) -> int:
        return page.page_number
    
    @staticmethod
    def get_pages(pdf: pdfplumber.PDF) -> int:
        return pdf.pages


class FitzParser(PDFParser):
    ''' fitz解析器实现 '''
    @staticmethod
    def open(pdf_path: str):
        return fitz.open(pdf_path)
    
    @staticmethod
    def get_page_text(page: fitz.Page) -> str:
        return page.get_text() or ""
    
    @staticmethod
    def get_page_number(page: fitz.Page) -> int:
        return page.number
    
    @staticmethod
    def get_pages(pdf: Any) -> int:
        return pdf
    

def get_parser(method: Literal['pdfplumber', 'fitz'] = 'pdfplumber'):
    if method == 'pdfplumber':
        return PDFPlumberParser
    elif method == 'fitz':
        return FitzParser
    raise ValueError(f"Unsupported method: {method}. Choose 'pdfplumber' or 'fitz'.")


def handle_cache(cache_dir: str, page_num: int, text: str) -> str:
    if not cache_dir:
        return text
    
    os.makedirs(cache_dir, exist_ok=True)
    file_path = os.path.join(cache_dir, f"page_{page_num}.txt")
    
    if os.path.exists(file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return text


def extract_pdf_single_process(
    pdf_path: str,
    cache_dir: str = None,
    page_start: int = 0,
    page_end: int = -1,
    method: Literal['pdfplumber', 'fitz'] = 'pdfplumber'
) -> List[str]:
    """
    单进程解析PDF文件
    
    :param pdf_path: PDF文件路径
    :param method: PDF解析方法
    :param cache_dir: 缓存目录路径
    :param page_start: 开始页码（从0开始）
    :param page_end: 结束页码（-1表示最后一页）
    :return: 按页码顺序排列的文本列表
    """
    all_text_list = []
    parser = get_parser(method)
    with parser.open(pdf_path) as pdf:
        total_pages = len(parser.get_pages(pdf))
        page_end = total_pages if page_end == -1 else min(page_end, total_pages)
        
        for page in tqdm(parser.get_pages(pdf)[page_start:page_end], desc='PDF extracting', ncols=80):
            page_num = parser.get_page_number(page)
            text = parser.get_page_text(page)
            cached_text = handle_cache(cache_dir, page_num, text)
            all_text_list.append(cached_text)
    
    return all_text_list


def _multiprocess_worker(args: Tuple[str, PDFParser, str, List[int]]) -> List[Tuple[int, str]]:
    ''' 多进程工作函数'''
    pdf_path, parser, cache_dir, page_indices = args
    results = []
    
    try:
        with parser.open(pdf_path) as pdf:
            for idx in tqdm(page_indices, desc='PDF extracting', ncols=80):
                if idx >= len(parser.get_pages(pdf)):
                    continue
                page = parser.get_pages(pdf)[idx]
                page_num = parser.get_page_number(page)
                text = parser.get_page_text(page)
                cached_text = handle_cache(cache_dir, page_num, text)
                results.append((page_num, cached_text))
    except Exception as e:
        log_error(f"processing pages {page_indices[0]}-{page_indices[-1]}: {e}")
    return results


def extract_pdf_multiprocess(
    pdf_path: str,
    cache_dir: str = None,
    page_start: int = 0,
    page_end: int = -1,
    num_processes: int = 4,
    method: Literal['pdfplumber', 'fitz'] = 'pdfplumber'
) -> List[str]:
    """
    多进程解析PDF文件
    
    :param pdf_path: PDF文件路径
    :param parser: PDF解析器实例
    :param cache_dir: 缓存目录路径
    :param page_start: 开始页码（从0开始）
    :param page_end: 结束页码（-1表示最后一页）
    :param num_processes: 进程数
    :return: 按页码顺序排列的文本列表
    """
    if cache_dir:
        os.makedirs(cache_dir, exist_ok=True)
    
    parser = get_parser(method)
    with parser.open(pdf_path) as pdf:
        total_pages = len(parser.get_pages(pdf))
        page_end = total_pages if page_end == -1 else min(page_end, total_pages)
    
    if page_start >= page_end:
        return []
    
    num_pages = page_end - page_start
    num_processes = max(min(num_processes, multiprocessing.cpu_count(), num_pages), 1)
    
    # 准备页码列表
    page_indices = list(range(page_start, page_end))
    # 分割页码列表
    chunk_size = (num_pages + num_processes - 1) // num_processes
    chunks = [page_indices[i:i+chunk_size] for i in range(0, num_pages, chunk_size)]
    
    # 多进程处理
    with multiprocessing.Pool(num_processes) as pool:
        args_list = [(pdf_path, parser, cache_dir, chunk) for chunk in chunks]
        results = list(pool.imap_unordered(_multiprocess_worker, args_list))
    
    # 整理结果
    all_text_list = [item for sublist in results for item in sublist]
    all_text_list.sort(key=lambda x: x[0])
    return [item[1] for item in all_text_list]


def _multithread_worker(parser:PDFParser, pdf_path:str, cache_dir:str, page_indices:List[int], results:List[str], lock:threading.Lock):
    ''' 多线程工作函数'''
    with parser.open(pdf_path) as pdf:
        for idx in tqdm(page_indices, desc='PDF extracting', ncols=80):
            if idx >= len(parser.get_pages(pdf)):
                continue
            page = parser.get_pages(pdf)[idx]
            page_num = parser.get_page_number(page)
            text = parser.get_page_text(page)
            cached_text = handle_cache(cache_dir, page_num, text)
            # 线程安全更新结果
            with lock:
                results[idx] = cached_text


def extract_pdf_multithread(
    pdf_path: str,
    cache_dir: str = None,
    page_start: int = 0,
    page_end: int = -1,
    num_threads: int = 4,
    method: Literal['pdfplumber', 'fitz'] = 'pdfplumber'
) -> List[str]:
    """
    多线程解析PDF文件
    
    :param pdf_path: PDF文件路径
    :param parser: PDF解析器实例
    :param cache_dir: 缓存目录路径
    :param page_start: 开始页码（从0开始）
    :param page_end: 结束页码（-1表示最后一页）
    :param num_threads: 线程数
    :return: 按页码顺序排列的文本列表
    """
    if cache_dir:
        os.makedirs(cache_dir, exist_ok=True)
    
    parser = get_parser(method)
    with parser.open(pdf_path) as pdf:
        total_pages = len(parser.get_pages(pdf))
        page_end = total_pages if page_end == -1 else min(page_end, total_pages)
    
    if page_start >= page_end:
        return []
    
    num_pages = page_end - page_start
    num_threads = max(min(num_threads, num_pages), 1)
    
    # 准备页码列表
    page_indices = list(range(page_start, page_end))
    results = [None] * num_pages
    lock = threading.Lock()
    sublist_len = (num_pages + num_threads - 1) // num_threads
    
    # 创建并启动线程
    threads = []
    for i in range(num_threads):
        start = i * sublist_len
        end = min(start + sublist_len, num_pages)
        thread = threading.Thread(
            target=_multithread_worker,
            args=(parser, pdf_path, cache_dir, page_indices[start:end], results, lock)
        )
        thread.start()
        threads.append(thread)
    
    # 等待所有线程完成
    for thread in threads:
        thread.join()
    
    return results


def extract_pdf_with_pdfplumber_coordinates(pdf_path, save_dir=None, page_start=0, page_end=-1, resolution=72):
    ''' 解析pdf并附带坐标信息（暂未使用） '''
    pdf = pdfplumber.open(pdf_path)

    page_end = len(pdf.pages) if page_end==-1 else page_end
    all_text_list = []
    for page in tqdm(pdf.pages[page_start:page_end], desc='PDF extracting with coordinates', ncols=80):
        file_path = os.path.join(save_dir, str(page.page_number)) + '_coords.jsonl' if save_dir is not None else ''
        if os.path.exists(file_path):
            # 文件存在
            text = []
            with open(file_path, 'r', encoding='utf-8') as f:
                for line in f.readlines():
                    text.append(json.loads(line))
        else:
            # 文件不存在
            text_elements = page.extract_words()  # 提取单词及其坐标
            text = [{'height': page.height, 'width': page.width, 'text': '', 'coords': [0, 0, 0, 0]}]
            for element in text_elements:
                text_tmp = element.get('text', '')  # 文本内容
                x1, y1, x2, y2 = element['x0'], element['top'], element['x1'], element['bottom']  # 坐标

                # 72 DPI 是 PDF 的默认分辨率
                x1 = int(x1 * resolution / 72)
                y1 = int(y1 * resolution / 72)
                x2 = int(x2 * resolution / 72)
                y2 = int(y2 * resolution / 72)

                text.append({'text': text_tmp, 'coords': [x1, y1, x2, y2]})
            if file_path != '':
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.writelines([json.dumps(i, ensure_ascii=False)+'\n' for i in text])

        all_text_list.append(text)
    return all_text_list


def dataframe_to_markdown(df):
    """
    将 Pandas DataFrame 转换为 Markdown 表格格式的字符串。
    
    参数:
    df (pd.DataFrame): 要转换的 DataFrame。
    
    返回:
    str: Markdown 表格格式的字符串。
    """
    # 获取列名
    columns = df.columns.tolist()
    
    # 构建表头
    header = "| " + " | ".join(columns) + " |"
    
    # 构建表头分隔线
    separator = "| " + " | ".join(["---"] * len(columns)) + " |"
    
    # 构建行数据
    rows = []
    for _, row in df.iterrows():
        values = []
        for i in row.values:
            if i is None or (isinstance(i, (int, float)) and np.isnan(i)):
                values.append('')
            elif isinstance(i, str):
                values.append(i.replace('\n', ''))
            else:
                values.append(i)
        row_data = "| " + " | ".join(map(str, values)) + " |"
        rows.append(row_data)
    
    # 合并所有部分
    markdown_table = "\n".join([header, separator] + rows)
    
    return markdown_table


def dataframe2text(df) -> str:
    '''dataframe转文本'''
    text_data = ", ".join(df.columns) + "\n"
    
    # 添加行数据
    for index, row in df.iterrows():
        row_data = ", ".join(map(str, row.values))
        text_data += row_data + "\n"
    
    text_data += "\n"  # 添加空行分隔不同的 sheet
    return text_data


def extract_txt(file_path:str) -> str:
    '''读取文本文件'''
    with open(file_path, 'r', encoding='utf-8') as f:
        data = f.read()
    return data


def extract_excel(file_path:str, df_format:Literal['markdown', 'text', 'dict', None]=None) -> list[str]:
    '''读取excel文件'''
    
    def is_row_empty(row):
        '''行是否为空'''
        return row.isnull().all()

    def find_first_non_empty_row_index(df):
        """找到第一个非空行的索引。"""
        for index, row in df.iterrows():
            if not is_row_empty(row):
                return index
        return None  # 如果所有行都是空的，返回 None
    
    import pandas as pd
    xls = pd.ExcelFile(file_path)
    sheet_names = xls.sheet_names
    
    all_text_data = []
    
    for sheet_name in sheet_names:
        # 读取每个 sheet 到 DataFrame
        df = pd.read_excel(xls, sheet_name=sheet_name)
        
        # 将 DataFrame 转换为文本格式
        text_data = f"Sheet Name: {sheet_name}\n\n"
        
        # 添加列名
        if all([re.search('Unnamed: [0-9]+', i) for i in df.columns]):
            # 空header
            header_index = find_first_non_empty_row_index(df)
            df.columns = df.iloc[header_index]
            df = df.drop(header_index, axis=0)

        if df_format == 'markdown':
            text_data += dataframe_to_markdown(df)
        elif df_format == 'text':
            text_data += dataframe2text(df)
        elif df_format == 'dict':
            text_data += "\n".join([json.dumps(i, ensure_ascii=False) for i in df.to_dict('records')])
        else:
            text_data = df
        all_text_data.append(text_data)
    
    return all_text_data

 
def extract_pptx(file_path) -> List[str]:
    '''从pptx中解析文档'''
    from pptx import Presentation
    presentation = Presentation(file_path)
    
    # 用于存储所有幻灯片的文本
    all_text = []
    
    # 遍历每个幻灯片
    for slide in presentation.slides:
        slide_text = []
        # 遍历每个形状
        for shape in slide.shapes:
            # 检查形状是否有文本框
            if shape.has_text_frame:
                # 提取文本框中的段落
                for paragraph in shape.text_frame.paragraphs:
                    slide_text.append(paragraph.text)
        # 将当前幻灯片的文本添加到总文本中
        all_text.append("\n".join(slide_text))
    
    return all_text